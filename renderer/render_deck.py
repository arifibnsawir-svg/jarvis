"""
render_deck.py v2 -- engine PRESENTASI editable (.pptx) ber-DESIGN-SYSTEM kuat.
Filosofi: "Jarvis berkreasi tanpa template" -- selera desain di-encode di kode (mirip CSS book),
bukan ngandelin file template. Output editable-native (bukan gambar), 16:9.

spec JSON -> render deterministik -> artifact_facts (anti-halu, buat validasi guard).
Dep: pip install python-pptx

Layout: cover | section | bullets | two_col | big_stat | quote | timeline | image | closing
Theme: bisa di-override via spec["theme"]. Preset: spec["preset"] in {academic, business, dark}.
"""
from __future__ import annotations
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- PRESET DESIGN TOKENS ----------
PRESETS = {
    "academic": {  # clean, scholarly, restrained
        "bg": "FFFFFF", "ink": "1A2233", "muted": "6B7689", "faint": "EEF1F6",
        "accent": "2F4B8E", "accent2": "1B2C54", "band_tx": "FFFFFF",
        "head_font": "Georgia", "body_font": "Calibri",
    },
    "business": {  # modern, confident
        "bg": "FFFFFF", "ink": "14171F", "muted": "667085", "faint": "F2F4F8",
        "accent": "2563EB", "accent2": "1E3A8A", "band_tx": "FFFFFF",
        "head_font": "Calibri", "body_font": "Calibri",
    },
    "dark": {
        "bg": "0B0D12", "ink": "ECEFF4", "muted": "9AA4B2", "faint": "1A1F29",
        "accent": "E08A4B", "accent2": "C25E2A", "band_tx": "0B0D12",
        "head_font": "Calibri", "body_font": "Calibri",
    },
}


def _rgb(h): return RGBColor.from_string(h.lstrip("#").upper())


def render_deck(spec: dict, out_path: str) -> dict:
    preset = PRESETS.get(spec.get("preset", "academic"), PRESETS["academic"])
    th = {**preset, **(spec.get("theme") or {})}
    BG, INK, MUTED, FAINT = th["bg"], th["ink"], th["muted"], th["faint"]
    ACC, ACC2, BANDTX = th["accent"], th["accent2"], th["band_tx"]
    HF, BF = th["head_font"], th["body_font"]
    footer = spec.get("footer", "")

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    ML = Inches(0.92); CW = W - 2 * ML

    def bg(slide, color):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = _rgb(color)

    def rect(slide, l, t, w, h, color):
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(color)
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
        x = slide.shapes.add_textbox(l, t, w, h); tf = x.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor; return tf

    def run(p, text, size, color, font, bold=False, italic=False, spacing=None):
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.name = font; r.font.color.rgb = _rgb(color)
        return r

    def eyebrow(slide, text, top=Inches(0.62)):
        tf = tb(slide, ML, top, CW, Inches(0.4))
        run(tf.paragraphs[0], text.upper(), 12, ACC, BF, bold=True)

    def title(slide, text, top=Inches(0.95)):
        tf = tb(slide, ML, top, CW, Inches(1.0))
        run(tf.paragraphs[0], text, 27, ACC2, HF, bold=True)
        rect(slide, ML, top + Inches(0.92), Inches(1.5), Emu(28575), ACC)

    def footer_bar(slide, idx):
        rect(slide, ML, H - Inches(0.6), CW, Emu(9525), FAINT)
        f1 = tb(slide, ML, H - Inches(0.54), CW * 0.75, Inches(0.35))
        run(f1.paragraphs[0], footer, 9.5, MUTED, BF)
        f2 = tb(slide, W - ML - Inches(1.0), H - Inches(0.54), Inches(1.0), Inches(0.35))
        f2.paragraphs[0].alignment = PP_ALIGN.RIGHT
        run(f2.paragraphs[0], f"{idx:02d}", 9.5, MUTED, BF)

    def bullets_block(slide, items, top, height, size=17, gap=10):
        body = tb(slide, ML, top, CW, height)
        for j, b in enumerate(items):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            run(p, "\u2014  ", size, ACC, BF, bold=True)  # em-dash style marker (akan tetap; visual, bukan prosa)
            if isinstance(b, str) and ": " in b and len(b.split(": ", 1)[0]) <= 30:
                lead, rest = b.split(": ", 1)
                run(p, lead, size, INK, BF, bold=True); run(p, " " + rest, size, INK, BF)
            else:
                run(p, str(b), size, INK, BF)
            p.space_after = Pt(gap); p.line_spacing = 1.18

    slides = spec.get("slides", []); page = 0; image_paths = []

    for s in slides:
        layout = s.get("layout", "bullets")
        sl = prs.slides.add_slide(blank); bg(sl, BG)

        if layout == "cover":
            rect(sl, 0, 0, Inches(0.45), H, ACC)
            if s.get("eyebrow"): 
                t = tb(sl, ML, Inches(2.0), CW, Inches(0.5)); run(t.paragraphs[0], s["eyebrow"].upper(), 13, ACC, BF, bold=True)
            t = tb(sl, ML, Inches(2.5), CW, Inches(2.2))
            run(t.paragraphs[0], s.get("title", ""), 42, INK, HF, bold=True)
            rect(sl, ML, Inches(4.05), Inches(2.2), Emu(38100), ACC)
            if s.get("subtitle"):
                st = tb(sl, ML, Inches(4.3), CW, Inches(1.2)); run(st.paragraphs[0], s["subtitle"], 19, MUTED, BF)
            if s.get("meta"):
                mt = tb(sl, ML, H - Inches(0.9), CW, Inches(0.5)); run(mt.paragraphs[0], "   |   ".join(s["meta"]), 11, MUTED, BF)

        elif layout == "section":
            rect(sl, 0, 0, W, H, ACC2)
            num = s.get("number")
            if num:
                nt = tb(sl, ML, Inches(2.2), CW, Inches(1.4)); run(nt.paragraphs[0], str(num), 64, ACC, HF, bold=True)
            t = tb(sl, ML, Inches(3.5), CW, Inches(1.4), MSO_ANCHOR.TOP)
            run(t.paragraphs[0], s.get("title", ""), 34, BANDTX, HF, bold=True)
            if s.get("subtitle"):
                p = t.add_paragraph(); run(p, s["subtitle"], 16, FAINT, BF)

        elif layout == "bullets":
            if s.get("eyebrow"): eyebrow(sl, s["eyebrow"])
            title(sl, s.get("title", ""), top=Inches(0.95) if s.get("eyebrow") else Inches(0.8))
            bullets_block(sl, s.get("bullets", []), Inches(2.2), H - Inches(3.0))
            page += 1; footer_bar(sl, page)

        elif layout == "two_col":
            title(sl, s.get("title", ""), top=Inches(0.8))
            colw = (CW - Inches(0.7)) / 2
            for ci, key in enumerate(("left", "right")):
                col = s.get(key) or {}; cl = ML + ci * (colw + Inches(0.7))
                if ci == 1: rect(sl, cl - Inches(0.35), Inches(2.15), Emu(9525), H - Inches(3.0), FAINT)
                if col.get("heading"):
                    ht = tb(sl, cl, Inches(2.05), colw, Inches(0.6)); run(ht.paragraphs[0], col["heading"], 18, ACC2, HF, bold=True)
                cb = tb(sl, cl, Inches(2.75), colw, H - Inches(3.6))
                for j, b in enumerate(col.get("bullets", [])):
                    p = cb.paragraphs[0] if j == 0 else cb.add_paragraph()
                    run(p, "\u2014  ", 15, ACC, BF, bold=True); run(p, str(b), 15, INK, BF)
                    p.space_after = Pt(9); p.line_spacing = 1.15
            page += 1; footer_bar(sl, page)

        elif layout == "big_stat":
            if s.get("eyebrow"): eyebrow(sl, s["eyebrow"])
            title(sl, s.get("title", ""), top=Inches(0.95) if s.get("eyebrow") else Inches(0.8))
            stats = s.get("stats", [])[:3]; n = len(stats) or 1
            gapw = Inches(0.5); cardw = (CW - gapw * (n - 1)) / n
            for i, stt in enumerate(stats):
                cx = ML + i * (cardw + gapw)
                rect(sl, cx, Inches(2.5), cardw, Inches(2.6), FAINT)
                vt = tb(sl, cx + Inches(0.25), Inches(2.75), cardw - Inches(0.5), Inches(1.2))
                run(vt.paragraphs[0], str(stt.get("value", "")), 44, ACC, HF, bold=True)
                lt = tb(sl, cx + Inches(0.25), Inches(3.95), cardw - Inches(0.5), Inches(1.0))
                run(lt.paragraphs[0], str(stt.get("label", "")), 14, INK, BF)
                if stt.get("note"):
                    nt = lt.add_paragraph(); run(nt, stt["note"], 10, MUTED, BF, italic=True)
            page += 1; footer_bar(sl, page)

        elif layout == "quote":
            rect(sl, ML, Inches(2.3), Inches(0.12), Inches(2.6), ACC)
            qt = tb(sl, ML + Inches(0.5), Inches(2.4), CW - Inches(0.5), Inches(2.4), MSO_ANCHOR.MIDDLE)
            run(qt.paragraphs[0], s.get("quote", ""), 28, INK, HF, italic=True)
            if s.get("attribution"):
                p = qt.add_paragraph(); run(p, "— " + s["attribution"], 15, MUTED, BF)
            page += 1; footer_bar(sl, page)

        elif layout == "timeline":
            title(sl, s.get("title", ""), top=Inches(0.8))
            phases = s.get("phases", []); n = len(phases) or 1
            gapw = Inches(0.4); cardw = (CW - gapw * (n - 1)) / n
            for i, ph in enumerate(phases):
                cx = ML + i * (cardw + gapw)
                rect(sl, cx, Inches(2.3), cardw, Emu(60000), ACC)  # bar atas
                rect(sl, cx, Inches(2.3), Inches(0.28), Inches(0.28), ACC2)  # node
                pt = tb(sl, cx, Inches(2.55), cardw, Inches(0.5))
                run(pt.paragraphs[0], str(ph.get("when", "")), 13, ACC, BF, bold=True)
                nt = tb(sl, cx, Inches(3.0), cardw, Inches(0.55))
                run(nt.paragraphs[0], str(ph.get("title", "")), 15, ACC2, HF, bold=True)
                dt = tb(sl, cx, Inches(3.6), cardw, H - Inches(4.4))
                run(dt.paragraphs[0], str(ph.get("desc", "")), 12.5, INK, BF); dt.paragraphs[0].line_spacing = 1.15
            page += 1; footer_bar(sl, page)

        elif layout == "image":
            title(sl, s.get("title", ""), top=Inches(0.8))
            ip = s.get("image_path", "")
            if not ip or not os.path.exists(ip):
                raise FileNotFoundError(f"image_path tidak ada (anti-halu): {ip!r}")
            image_paths.append(ip); sl.shapes.add_picture(ip, ML, Inches(2.2), width=CW)
            page += 1; footer_bar(sl, page)

        elif layout == "closing":
            rect(sl, 0, 0, W, H, ACC2)
            rect(sl, ML, Inches(2.5), Inches(2.2), Emu(38100), ACC)
            t = tb(sl, ML, Inches(2.8), CW, Inches(1.6))
            run(t.paragraphs[0], s.get("title", "Terima Kasih"), 40, BANDTX, HF, bold=True)
            if s.get("subtitle"):
                st = tb(sl, ML, Inches(4.2), CW, Inches(1.2)); run(st.paragraphs[0], s["subtitle"], 17, FAINT, BF)
            for i, c in enumerate(s.get("contacts", [])):
                ct = tb(sl, ML, H - Inches(1.6) + Inches(0.4) * i, CW, Inches(0.4))
                run(ct.paragraphs[0], c, 13, FAINT, BF)
        else:
            raise ValueError(f"layout tidak dikenal: {layout}")

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)

    def _texts():
        out = []
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.has_text_frame: out.append(sh.text_frame.text)
        return " ".join(out)

    return {
        "artifact": out_path, "slide_count": len(slides), "image_paths": image_paths,
        "ratio": "16:9", "preset": spec.get("preset", "academic"), "theme_accent": "#" + ACC,
        "word_count": len(_texts().split()), "file_size": os.path.getsize(out_path),
    }


if __name__ == "__main__":
    import json, sys
    if len(sys.argv) < 3:
        print("usage: render_deck.py spec.json out.pptx"); sys.exit(1)
    print(json.dumps(render_deck(json.load(open(sys.argv[1])), sys.argv[2]), ensure_ascii=False, indent=2))

"""
PPTX renderer — bangun carousel/microblog dari spec JSON (IG/LinkedIn).
Deterministik: spec sama -> file sama. Return artifact_facts utk GUARDIAN.

Dep: pip install python-pptx
"""
from __future__ import annotations

import os
from spec_schema import validate_spec, extract_text

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DEFAULT_THEME = {"bg": "#0E1116", "fg": "#E6E6E6", "accent": "#4F8CFF", "font": "Calibri"}


def _rgb(hexs: str) -> RGBColor:
    return RGBColor.from_string(hexs.lstrip("#").upper())


def _fill_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    return tb.text_frame


def _set(run, text, size, color, font, bold=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def render_pptx(spec: dict, out_path: str) -> dict:
    spec = validate_spec(spec)
    theme = {**DEFAULT_THEME, **spec.get("theme", {})}
    bg, fg, accent, font = _rgb(theme["bg"]), _rgb(theme["fg"]), _rgb(theme["accent"]), theme["font"]

    prs = Presentation()
    prs.slide_width = Inches(7.5)      # square 1:1 (carousel)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height
    M = Inches(0.6)                    # margin
    image_paths = []

    for s in spec["slides"]:
        slide = prs.slides.add_slide(blank)
        _fill_bg(slide, bg)
        layout = s["layout"]

        if layout == "title":
            tf = _textbox(slide, M, Inches(2.6), W - 2 * M, Inches(2.3), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            _set(p.add_run(), s["title"], 40, fg, font, bold=True)
            if s.get("subtitle"):
                p2 = tf.add_paragraph(); _set(p2.add_run(), s["subtitle"], 20, accent, font)

        elif layout == "bullets":
            tf = _textbox(slide, M, M, W - 2 * M, Inches(1.4))
            _set(tf.paragraphs[0].add_run(), s["title"], 30, accent, font, bold=True)
            body = _textbox(slide, M, Inches(2.2), W - 2 * M, H - Inches(2.8))
            for j, b in enumerate(s["bullets"]):
                p = body.paragraphs[0] if j == 0 else body.add_paragraph()
                _set(p.add_run(), f"•  {b}", 20, fg, font)
                p.space_after = Pt(10)

        elif layout == "body":
            tf = _textbox(slide, M, M, W - 2 * M, Inches(1.4))
            _set(tf.paragraphs[0].add_run(), s["title"], 30, accent, font, bold=True)
            body = _textbox(slide, M, Inches(2.2), W - 2 * M, H - Inches(2.8))
            _set(body.paragraphs[0].add_run(), s["body"], 20, fg, font)

        elif layout == "quote":
            tf = _textbox(slide, M, Inches(2.2), W - 2 * M, Inches(3), MSO_ANCHOR.MIDDLE)
            _set(tf.paragraphs[0].add_run(), f"\u201c{s['quote']}\u201d", 28, fg, font, bold=True)
            if s.get("attribution"):
                p = tf.add_paragraph(); _set(p.add_run(), f"— {s['attribution']}", 18, accent, font)

        elif layout == "image":
            if not os.path.exists(s["image_path"]):
                raise FileNotFoundError(f"image_path tidak ada (anti-halu): {s['image_path']}")
            image_paths.append(s["image_path"])
            slide.shapes.add_picture(s["image_path"], M, M, width=W - 2 * M)
            if s.get("title"):
                tf = _textbox(slide, M, H - Inches(1.2), W - 2 * M, Inches(0.9))
                _set(tf.paragraphs[0].add_run(), s["title"], 22, fg, font, bold=True)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)

    facts = {
        "artifact": out_path,
        "slide_count": len(spec["slides"]),
        "image_paths": image_paths,
        "word_count": len(extract_text(spec).split()),
        "file_size": os.path.getsize(out_path),
    }
    return facts
